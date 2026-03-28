from __future__ import annotations

import sys
from pathlib import Path

SITE_PACKAGES = Path(
    r"C:\Users\KIIT0001\AppData\Local\Programs\Python\Python313\Lib\site-packages"
)
if str(SITE_PACKAGES) not in sys.path:
    sys.path.append(str(SITE_PACKAGES))

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROOT = Path(r"C:\Users\KIIT0001\Desktop\New folder\subfolder")
SCREENSHOT_DIR = ROOT / "presentation_assets" / "screenshots"

INPUT_DECK = Path(r"C:\Users\KIIT0001\Downloads\Social Engineering Benchmark Presentation Deck.pptx")
OUTPUT_DECK = Path(r"C:\Users\KIIT0001\Downloads\Social Engineering Benchmark Presentation Deck_fixed.pptx")

ENV_ORDER = [
    "e7-ninite",
    "e8-chase",
    "e9-expedia",
    "e10-mychart",
    "e11-instagram",
    "e12-dhl",
    "e13-irs",
    "e14-oracle",
    "e15-metamask",
    "e16-microsoft",
    "e17-amazon",
    "e17b-servicenow",
    "e18-linkedin",
    "e19-apple",
    "e20-survey",
    "e21-paypal",
    "e22-zoom",
    "e23-netflix",
    "e24-uber",
    "e25-airbnb",
]


def replace_picture(shape, image_path: Path) -> None:
    shapes = shape._parent
    old_element = shape._element
    parent = old_element.getparent()
    insert_at = parent.index(old_element)

    new_shape = shapes.add_picture(
        str(image_path),
        shape.left,
        shape.top,
        shape.width,
        shape.height,
    )
    new_shape.crop_left = shape.crop_left
    new_shape.crop_right = shape.crop_right
    new_shape.crop_top = shape.crop_top
    new_shape.crop_bottom = shape.crop_bottom

    new_element = new_shape._element
    parent.remove(new_element)
    parent.insert(insert_at, new_element)
    parent.remove(old_element)


def main() -> None:
    prs = Presentation(INPUT_DECK)

    for offset, env_name in enumerate(ENV_ORDER, start=4):
        slide = prs.slides[offset]
        pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
        if not pictures:
            raise RuntimeError(f"No picture found on slide {offset + 1} for {env_name}")
        replace_picture(pictures[0], SCREENSHOT_DIR / f"{env_name}.png")

    prs.save(OUTPUT_DECK)
    print(OUTPUT_DECK)


if __name__ == "__main__":
    main()
