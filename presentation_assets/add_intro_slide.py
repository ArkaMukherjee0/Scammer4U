from __future__ import annotations

import sys
from pathlib import Path

SITE_PACKAGES = Path(
    r"C:\Users\KIIT0001\AppData\Local\Programs\Python\Python313\Lib\site-packages"
)
if str(SITE_PACKAGES) not in sys.path:
    sys.path.append(str(SITE_PACKAGES))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


INPUT_DECK = Path(
    r"C:\Users\KIIT0001\Desktop\New folder\subfolder\presentation_assets\output\Scammer4U_Benchmark_Deck.pptx"
)
OUTPUT_DECK = Path(
    r"C:\Users\KIIT0001\Desktop\New folder\subfolder\presentation_assets\output\Scammer4U_Benchmark_Deck_named.pptx"
)


def move_last_slide_to_front(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst
    new_slide_id = sld_id_lst[-1]
    sld_id_lst.remove(new_slide_id)
    sld_id_lst.insert(0, new_slide_id)


def main() -> None:
    prs = Presentation(INPUT_DECK)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        prs.slide_width,
        prs.slide_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(18, 24, 38)
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.53), Inches(2.0))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_run = title_para.add_run()
    title_run.text = "Anyash Prasad"
    title_run.font.name = "Aptos Display"
    title_run.font.size = Pt(30)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(245, 247, 250)

    subtitle_box = slide.shapes.add_textbox(Inches(2.35), Inches(3.05), Inches(8.65), Inches(1.0))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.clear()
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_run = subtitle_para.add_run()
    subtitle_run.text = "Roll No. - 23051977"
    subtitle_run.font.name = "Aptos"
    subtitle_run.font.size = Pt(20)
    subtitle_run.font.color.rgb = RGBColor(192, 200, 214)

    footer_box = slide.shapes.add_textbox(Inches(3.15), Inches(4.65), Inches(7.0), Inches(0.6))
    footer_frame = footer_box.text_frame
    footer_frame.clear()
    footer_para = footer_frame.paragraphs[0]
    footer_para.alignment = PP_ALIGN.CENTER
    footer_run = footer_para.add_run()
    footer_run.text = "Social Engineering Benchmark Presentation"
    footer_run.font.name = "Aptos"
    footer_run.font.size = Pt(16)
    footer_run.font.color.rgb = RGBColor(133, 146, 166)

    move_last_slide_to_front(prs)
    prs.save(OUTPUT_DECK)
    print(OUTPUT_DECK)


if __name__ == "__main__":
    main()
